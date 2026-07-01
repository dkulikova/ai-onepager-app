from __future__ import annotations

import json
import os
import re
import sqlite3
import textwrap
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import quote_plus

import feedparser
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).parent


def resolve_db_path() -> Path:
    candidates = [BASE / "data" / "briefing_demo.db", BASE / "briefing_demo.db"]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return BASE / "data" / "briefing_demo.db"


def resolve_notes_dirs() -> List[Path]:
    candidates = [BASE / "notes", BASE]
    return [path for path in candidates if path.exists()]


DB_PATH = resolve_db_path()
st.set_page_config(page_title="BriefIQ", layout="wide")


# -----------------------------
# Data access layer
# -----------------------------
@st.cache_resource
def get_connection(db_path_str: str) -> sqlite3.Connection:
    return sqlite3.connect(db_path_str, check_same_thread=False)


@st.cache_data
def safe_load_table(table_name: str, db_path_str: str) -> pd.DataFrame:
    conn = get_connection(db_path_str)
    try:
        return pd.read_sql_query(f"SELECT * FROM {table_name}", conn)
    except Exception:
        return pd.DataFrame()


def read_markdown_notes() -> Dict[str, str]:
    notes: Dict[str, str] = {}
    for notes_dir in resolve_notes_dirs():
        for file in notes_dir.glob("*.md"):
            if file.name.lower() == "readme.md":
                continue
            if file.name in notes:
                continue
            notes[file.name] = file.read_text(encoding="utf-8", errors="ignore")
    return notes


def retrieve_relevant_notes(query: str, notes: Dict[str, str], max_notes: int = 3) -> List[str]:
    query_terms = {term.lower().strip() for term in re.split(r"\W+", query) if len(term) > 2}
    scored = []
    for name, text in notes.items():
        haystack = f"{name}\n{text}".lower()
        score = sum(1 for term in query_terms if term in haystack)
        if score > 0:
            scored.append((score, name, text))
    scored.sort(reverse=True, key=lambda x: x[0])
    return [f"Source note: {name}\n{text}" for _, name, text in scored[:max_notes]]


def table_for_company(table_name: str, company: pd.Series) -> pd.DataFrame:
    df = safe_load_table(table_name, str(DB_PATH))
    if df.empty or "company_id" not in df.columns or "company_id" not in company.index:
        return pd.DataFrame()
    company_id = company.get("company_id")
    return df.loc[df["company_id"].astype(str) == str(company_id)].copy()


def company_intelligence(company: pd.Series) -> Dict[str, pd.DataFrame]:
    return {
        "leadership": table_for_company("leadership", company),
        "products": table_for_company("products", company),
        "financials": table_for_company("financials", company),
        "milestones": table_for_company("milestones", company),
    }


def format_df_records(title: str, df: pd.DataFrame) -> str:
    if df is None or df.empty:
        return f"{title}: No internal records found."
    lines = [title]
    for _, row in df.iterrows():
        values = []
        for col, val in row.items():
            if col.endswith("_id") or col == "company_id":
                continue
            values.append(f"{col}: {val}")
        lines.append("- " + "; ".join(values))
    return "\n".join(lines)


@st.cache_data(ttl=1800, show_spinner=False)
def get_latest_company_articles(company_name: str, max_articles: int = 5) -> List[Dict[str, str]]:
    """Return the first Google News RSS results for the typed company name.

    This deliberately uses the company name as the lookup term and returns the
    first five non-empty RSS entries. The LLM and human reviewer then decide
    whether those items are relevant enough to include in the profile.
    """
    lookup_term = str(company_name or "").strip()
    query = quote_plus(lookup_term)
    rss_url = f"https://news.google.com/rss/search?q={query}&hl=en-GB&gl=GB&ceid=GB:en"

    try:
        feed = feedparser.parse(rss_url)
    except Exception as exc:
        return [{"title": "Latest news retrieval failed", "published": "", "summary": f"The RSS feed could not be read: {exc}", "link": rss_url, "source": "Google News RSS"}]

    articles: List[Dict[str, str]] = []
    for entry in feed.entries:
        title = re.sub(r"\s+", " ", entry.get("title", "")).strip()
        if not title:
            continue
        summary = re.sub(r"<.*?>", "", entry.get("summary", "")).strip()
        source = ""
        try:
            source = entry.get("source", {}).get("title", "") if isinstance(entry.get("source", {}), dict) else ""
        except Exception:
            source = ""
        articles.append({
            "title": title,
            "published": entry.get("published", ""),
            "summary": summary,
            "link": entry.get("link", ""),
            "source": source,
            "lookup_term": lookup_term,
        })
        if len(articles) >= max_articles:
            break

    if not articles:
        return [{"title": "No recent articles retrieved", "published": "", "summary": f"No RSS results were returned for lookup term: {lookup_term}.", "link": rss_url, "source": "Google News RSS"}]
    return articles


def format_articles_for_prompt(articles: List[Dict[str, str]]) -> str:
    if not articles:
        return "Latest news was excluded or no recent articles were retrieved."
    formatted = []
    for i, article in enumerate(articles, start=1):
        formatted.append(
            f"""
Article {i}
Title: {article.get('title', '')}
Published: {article.get('published', '')}
Summary: {article.get('summary', '')}
Link: {article.get('link', '')}
""".strip()
        )
    return "\n\n".join(formatted)


def format_article_for_slide(article: Dict[str, str]) -> str:
    """Format a raw RSS entry for the PowerPoint news panel without adding ellipses."""
    title = re.sub(r"\s+", " ", str(article.get("title", "")).strip())
    source = str(article.get("source", "")).strip()
    published = str(article.get("published", "")).strip()
    # Google News often appends the publisher after " - "; keep the title as returned
    # because the user explicitly wanted the first five RSS items for the company lookup term.
    parts = [title]
    if source and source.lower() not in title.lower():
        parts.append(source)
    if published:
        parts.append(published[:16])
    return " — ".join([p for p in parts if p])


# -----------------------------
# Matching / project context
# -----------------------------
def find_company_match(company_query: str, companies_df: pd.DataFrame) -> tuple[pd.Series, bool, str]:
    query = (company_query or "").strip() or "Unknown company"
    names = companies_df["company_name"].astype(str).tolist() if not companies_df.empty else []

    lower_map = {name.lower(): name for name in names}
    if query.lower() in lower_map:
        matched_name = lower_map[query.lower()]
        row = companies_df.loc[companies_df["company_name"] == matched_name].iloc[0]
        return row, True, f"Exact internal database match found: {matched_name}."

    for name in names:
        if query.lower() in name.lower() or name.lower() in query.lower():
            row = companies_df.loc[companies_df["company_name"] == name].iloc[0]
            return row, True, f"Likely internal database match found: {name}."

    scored = [(SequenceMatcher(None, query.lower(), name.lower()).ratio(), name) for name in names]
    scored.sort(reverse=True)
    if scored and scored[0][0] >= 0.72:
        matched_name = scored[0][1]
        row = companies_df.loc[companies_df["company_name"] == matched_name].iloc[0]
        return row, True, f"Fuzzy internal database match found: {matched_name} ({scored[0][0]:.0%} similarity)."

    synthetic = {col: "Not available" for col in companies_df.columns} if not companies_df.empty else {}
    synthetic.update({
        "company_id": "external_query",
        "company_name": query,
        "sector": "Not available from internal database",
        "hq_location": "Not available from internal database",
        "founded_year": "Not available",
        "company_type": "Not available",
        "employee_count": "Not available",
        "website": "Not available",
        "current_office_location": "Not available",
        "mission": "No internal mission record found.",
        "short_description": "No internal company profile found. Use user context and optional latest news cautiously.",
        "recent_news": "No internal database record found.",
        "expansion_signal": "No internal expansion signal found. Requires external verification.",
        "target_market": "Not available",
        "differentiation": "Not available",
        "relevance_score": "N/A",
    })
    return pd.Series(synthetic), False, "No confident internal database match found. The one-pager will use user context and optional latest news."


def build_project_context(brief_type: str, extra_notes: str) -> pd.Series:
    descriptions = {
        "Investment one-pager": "Create a concise investment-style profile covering mission, leadership, products, funding/commercial signals, risks and next diligence steps.",
        "Company brief": "Create a concise company profile covering mission, leadership, what the company does, key facts, recent signals and business relevance.",
    }
    return pd.Series({
        "project_id": "user_defined",
        "project_name": brief_type,
        "client": "User-defined / interview demo",
        "location": "Not specified",
        "target_sector": "To be inferred from the company and context",
        "brief_description": descriptions.get(brief_type, "Create a concise business one-pager."),
        "user_context": extra_notes or "No additional context provided.",
    })


# -----------------------------
# Source pack / prompts
# -----------------------------
def build_source_pack(
    selected_company: pd.Series,
    selected_project: pd.Series,
    extra_notes: str,
    retrieved_notes: List[str],
    latest_articles: List[Dict[str, str]],
    intel: Dict[str, pd.DataFrame],
    has_internal_match: bool = True,
    match_message: str = "",
    live_leadership_lookup: str = "",
    live_company_profile_lookup: str = "",
) -> str:
    company_context = "\n".join([f"- {col}: {selected_company[col]}" for col in selected_company.index])
    project_context = "\n".join([f"- {col}: {selected_project[col]}" for col in selected_project.index])
    note_context = "\n\n".join(retrieved_notes) if retrieved_notes else "No relevant local notes found."
    latest_news_context = format_articles_for_prompt(latest_articles)
    database_status = (
        "Internal database match found. Use internal structured records as the primary evidence."
        if has_internal_match
        else "NO INTERNAL DATABASE MATCH FOUND. The company was not found in the internal database. Fill only what can be supported from user context, local notes, latest news/RSS and very high-confidence general knowledge; add a clear verification warning at the beginning of the output."
    )

    return f"""
INTERNAL DATABASE MATCH STATUS
- has_internal_match: {has_internal_match}
- match_message: {match_message}
- instruction: {database_status}

STRUCTURED COMPANY DATA
{company_context}

EXECUTIVE LEADERSHIP DATA
{format_df_records('Leadership records', intel.get('leadership', pd.DataFrame()))}

LIVE COMPANY PROFILE WEB LOOKUP
{live_company_profile_lookup or 'No live company profile lookup was run or no live company profile results were returned.'}

LIVE LEADERSHIP WEB LOOKUP
{live_leadership_lookup or 'No live leadership lookup was run or no live leadership results were returned.'}

PRODUCTS / SERVICES DATA
{format_df_records('Product records', intel.get('products', pd.DataFrame()))}

FUNDING / COMMERCIAL SIGNALS DATA
{format_df_records('Financial and commercial records', intel.get('financials', pd.DataFrame()))}

MILESTONES DATA
{format_df_records('Milestone records', intel.get('milestones', pd.DataFrame()))}

PROJECT / USE CASE DATA
{project_context}

LOCAL MARKDOWN / OBSIDIAN NOTES
{note_context}

LATEST NEWS / RECENT EXTERNAL SIGNALS
{latest_news_context}

USER ADDED NOTES
{extra_notes or 'No additional notes provided.'}
""".strip()


def build_prompt(source_pack: str, brief_type: str) -> str:
    """Main generation prompt for the OpenAI call.

    This is the core prompt-engineering layer. It tells the model what role to play,
    what evidence to trust, how to treat private-company financial information, and
    exactly which structured fields to return for the PowerPoint template.
    """
    return f"""
You are a senior business analyst creating a one-page company intelligence brief for a non-technical senior audience.

Brief type: {brief_type}

Your job is to turn the source pack into a concise, evidence-grounded PowerPoint one-pager.
You are not writing marketing copy. You are creating a practical business briefing that a human can review, edit and use.

Core rules:
- First check the INTERNAL DATABASE MATCH STATUS in the source pack.
- If an internal database match is available, use the internal structured records as the primary evidence.
- If no internal database match is available, still produce a useful first draft by using the user context, local notes if any, latest news/RSS if included, and your high-confidence general knowledge.
- If no internal database match is available, include this exact warning in the `verification_banner` field: "Not found in the internal database — this profile uses external/contextual information and needs additional verification."
- Your own model knowledge may be older than today's date. Do not treat your pretrained knowledge as live data. For time-sensitive facts such as current global CEO, global executive leadership, valuation, funding, employee count and recent announcements, prefer internal database records, LIVE LEADERSHIP WEB LOOKUP, user-provided context and RSS/news items. If relying on general knowledge, explicitly mark the item "to verify".
- Do not fabricate or guess precise facts. Only include facts you believe are true and stable. Where you are not confident, write "to verify" or "not available".
- For leadership, return exactly 3 people. Prioritise: (1) current global CEO, (2) current global executive leadership such as CFO, COO, President, CTO, Chief Product Officer, Chief Scientist, Chair/Founder where materially relevant. Do not include regional leaders, board-only roles, retired executives or former CEOs unless the role is clearly current in the source pack.
- For leadership when no internal database record exists: use the LIVE LEADERSHIP WEB LOOKUP if available. If no live lookup is available, include the current global CEO and two current global executive leaders you believe are correct from high-confidence general knowledge, but mark each item "to verify". Do not include obviously historic or retired leaders.
- For milestones when no internal database record exists: populate the timeline with specific, recognisable company events such as founded year, major product launches, funding rounds, public listings, acquisitions or recent announcements from RSS/user context. Avoid generic milestones such as "growth continues". Mark uncertain dates or announcements as "to verify".
- If the source pack is incomplete, say what is missing or what needs verification. For companies not found in the database, use the live company profile web lookup to populate HQ, founded year, company type, sector and employee scale where available.
- Treat private-company funding, valuation and revenue references as funding/commercial signals, not as audited financial performance.
- Separate facts from interpretation. Do not make weak evidence sound definitive.
- Latest news/RSS headlines are external signals to verify, not proof on their own. Use the first five RSS results provided in the source pack as the raw latest-news feed. Do not replace them with older model knowledge. If one appears irrelevant, label it "verify relevance" rather than inventing a different item.
- Write for a smart business audience that may not be technical.
- Use clear, direct language and avoid AI hype.

Prioritise:
1. Mission, positioning and strategic focus
2. Executive leadership and why the leadership context matters
3. Products/services and what the company actually does
4. Funding/commercial signals, especially for private companies
5. Latest news only where it is relevant to the brief
6. Risks, caveats and questions for human review
7. Concrete next steps

Brief-type guidance:
- If the brief type is "Investment one-pager", emphasise market relevance, leadership, product surface, growth/funding signals, risks, and diligence next steps.
- If the brief type is "Company brief", emphasise what the company does, mission/positioning, leadership, key facts, recent signals, and business relevance.

SOURCE PACK:
{source_pack}

Return ONLY valid JSON. Do not include markdown, commentary, citations outside the JSON, or code fences.
Use exactly this schema. Write complete but compact sentences that fit a single PowerPoint slide. Do not use ellipses or unfinished text:
{{
  "headline": "short title for the profile",
  "verification_banner": "use an empty string if the company was found in the internal database; otherwise include the required not-in-database verification warning",
  "company_positioning": "2-3 complete short lines on mission, positioning and differentiation",
  "growth_direction": "2-3 complete short lines on growth direction and strategic momentum",
  "target_market": "2-3 complete short lines on target users, customers or market",
  "company_description": "two complete sentences, maximum 45 words total, suitable for a left-side profile panel",
  "company_snapshot": {{"hq": "global HQ or to verify", "founded": "founded year or to verify", "company_type": "private/public/subsidiary/status or to verify", "sector": "sector or to verify", "employees": "employee count/scale or to verify"}},
  "what_they_do": ["complete sentence, maximum 14 words", "complete sentence, maximum 14 words", "complete sentence, maximum 14 words"],
  "leadership": ["Current global CEO — role; brief relevance", "Current global executive leader — role; brief relevance", "Current global executive leader — role; brief relevance"],
  "key_facts": ["specific fact 1", "specific fact 2", "specific fact 3", "specific fact 4"],
  "funding_commercial_signals": ["complete sentence, maximum 16 words", "complete sentence, maximum 16 words", "complete sentence, maximum 16 words"],
  "latest_news_signals": ["article title 1 — publisher/source if available", "article title 2 — publisher/source if available", "article title 3 — publisher/source if available", "article title 4 — publisher/source if available", "article title 5 — publisher/source if available"],
  "risks": ["complete sentence, maximum 14 words", "complete sentence, maximum 14 words", "complete sentence, maximum 14 words"],
  "timeline": [
    {{"year": "2023", "text": "milestone or signal"}},
    {{"year": "2024", "text": "milestone or signal"}},
    {{"year": "2025", "text": "milestone or signal"}},
    {{"year": "2026", "text": "milestone, signal or to verify"}}
  ],
  "next_steps": ["complete sentence, maximum 14 words", "complete sentence, maximum 14 words", "complete sentence, maximum 14 words"]
}}

Field-specific instructions:
- Write complete sentences that fit the slide. Do not use ellipses, sentence fragments, trailing clauses, or intentionally cut-off wording. Prefer shorter complete sentences over long sentences.
- "company_positioning", "growth_direction" and "target_market": write one compact complete sentence for each, around 18-26 words. Do not insert line breaks or ellipses. The PowerPoint template will wrap the sentence into the original top-card text box.
- "verification_banner": if no internal database match is available, this must be the first visible message in the profile. If there is an internal match, leave it blank.
- "company_snapshot": populate the left-hand company snapshot fields. If the company is found in the database, use the internal database first. If it is not found, use the LIVE COMPANY PROFILE WEB LOOKUP, latest news/user context and high-confidence general knowledge. If a field is uncertain, write "to verify" rather than leaving it blank.
- "leadership": return exactly 3 current global leaders. Include the current global CEO as the first item wherever reasonably available. If internal leadership records exist, use them but limit to the three most relevant global executives. If not, use the LIVE LEADERSHIP WEB LOOKUP. If no live lookup is available, use high-confidence general knowledge and add "to verify" where appropriate. Do not include retired/former CEOs or regional leaders unless explicitly current and globally relevant.
- "funding_commercial_signals": for private companies, use language such as "reported", "funding signal", "commercial signal", or "to verify" where appropriate.
- "latest_news_signals": use the first five RSS results from the source pack as the feed. Return five article-style bullets wherever five RSS items are present. Use the article title and source/publisher where available. Do not leave empty rows. Do not silently replace RSS items with older model knowledge.
- "timeline": if internal milestones exist, use them. If no internal milestones exist, create a useful, chronological timeline from high-confidence founding year, product launches, funding/company announcements, public listings/acquisitions and relevant RSS/user-context signals. Each milestone must be a complete compact sentence of 8 words or fewer. Avoid generic wording such as "growth continues". Use "to verify" for uncertain items.
- "risks": include both information-quality risks and business risks where relevant.
- "next_steps": make these practical actions a human analyst or business team would take before using the one-pager.
""".strip()


def build_review_prompt(profile_json: str, source_pack: str) -> str:
    """Quality-control prompt for the second OpenAI call."""
    return f"""
You are reviewing an AI-generated company one-pager before it is used by a business team.

Compare the profile JSON against the source pack. Be strict and practical.

Return a concise review with four headings:

1. Claims to verify
- Flag any claim that appears unsupported, too broad, too confident, or dependent on a weak source.

2. Missing information
- Highlight any important missing details, especially leadership, mission, products, funding/commercial signals, financials, or recent news.

3. Risk notes
- Flag private-company financial uncertainty, duplicated/irrelevant RSS results, outdated data, or areas where human judgement is needed.

4. Suggested improvements for a non-technical audience
- Suggest how to make the one-pager clearer, more specific, or less jargon-heavy.

Rules:
- Do not rewrite the full one-pager.
- Do not introduce new facts.
- Be concise, direct and useful.
- Treat latest news as signals to verify, not definitive evidence.
- Check that the leadership section includes exactly 3 people, starting with the current global CEO where available, and does not include retired/former CEOs or regional-only leaders.
- If the company was not found in the internal database, check that the profile clearly starts with the not-in-database verification warning.

SOURCE PACK:
{source_pack}

PROFILE JSON:
{profile_json}
""".strip()


# -----------------------------
# LLM provider layer
# -----------------------------
def get_secret(name: str) -> Optional[str]:
    try:
        value = st.secrets.get(name)
        if value:
            return str(value)
    except Exception:
        pass
    return os.getenv(name)


def call_openai(prompt: str, model: str = "gpt-4.1-mini") -> str:
    from openai import OpenAI

    client = OpenAI(api_key=get_secret("OPENAI_API_KEY"))
    response = client.responses.create(model=model, input=prompt)
    return response.output_text






def lookup_current_company_profile(company_name: str) -> str:
    """Use OpenAI hosted web search to retrieve current company profile fields for external companies.

    This is used when the company is not in the internal database. It gives the
    main generation prompt fresher context for the left-hand snapshot fields,
    mission/positioning and milestone suggestions.
    """
    api_key = get_secret("OPENAI_API_KEY")
    if not api_key:
        return "No live company profile lookup was run because OPENAI_API_KEY is not configured."

    from openai import OpenAI
    client = OpenAI(api_key=api_key)
    prompt = f"""
Find current, high-confidence company profile information for: {company_name}.

Return concise plain text with these headings:
- Global HQ
- Founded year
- Company type / listing status
- Sector
- Employees / scale, if available
- Mission or company positioning
- Main products/services
- Recent milestones or announcements

Rules:
- Prioritise official company pages, investor relations pages, annual reports, company newsroom announcements, or highly reputable business/news sources.
- Use global company information, not regional subsidiaries, unless the company itself is regional.
- If a field is not confidently available, write "not available / to verify".
- Do not use stale leadership or retired executives.
- Keep the output concise; this will be used as source material for a one-slide briefing.
""".strip()

    try:
        response = client.responses.create(
            model="gpt-4.1-mini",
            tools=[{"type": "web_search"}],
            tool_choice="required",
            input=prompt,
        )
        return response.output_text.strip()
    except Exception as exc:
        return f"Live company profile lookup failed or is unavailable: {exc}"


def lookup_current_global_leadership(company_name: str) -> str:
    """Use OpenAI hosted web search to retrieve current global leadership for an external company.

    This is only used as additional source-pack context. The final profile prompt still
    has to decide what to include and mark uncertain items for verification.
    """
    api_key = get_secret("OPENAI_API_KEY")
    if not api_key:
        return "No live leadership lookup was run because OPENAI_API_KEY is not configured."

    from openai import OpenAI
    client = OpenAI(api_key=api_key)
    prompt = f"""
Find the current global CEO and current global executive leadership for: {company_name}.

Return ONLY concise plain text with exactly three bullets:
1. Current global CEO — role — source or verification note
2. Current global executive leader — role — source or verification note
3. Current global executive leader — role — source or verification note

Rules:
- Prioritise official company leadership pages, investor relations pages, company newsroom announcements, or highly reputable business/news sources.
- Use global company leadership only. Do not use regional country heads unless no global leadership information is available.
- Do not include retired leaders, former CEOs, former founders in non-current roles, or board-only figures unless clearly part of current executive leadership.
- If a current role is uncertain, write "to verify".
""".strip()

    try:
        response = client.responses.create(
            model="gpt-4.1-mini",
            tools=[{"type": "web_search"}],
            tool_choice="required",
            input=prompt,
        )
        return response.output_text.strip()
    except Exception as exc:
        return f"Live leadership lookup failed or is unavailable: {exc}"


def run_llm(prompt: str, provider: str = "OpenAI") -> str:
    """Runs the generation or review prompt using OpenAI.

    If no OpenAI key is available, return an empty string so the app can fall back
    to the deterministic demo output instead of crashing during an interview demo.
    """
    if provider == "OpenAI" and get_secret("OPENAI_API_KEY"):
        return call_openai(prompt)
    return ""


# -----------------------------
# Profile parsing / formatting
# -----------------------------
def extract_json(text: str) -> Optional[Dict[str, Any]]:
    try:
        return json.loads(text)
    except Exception:
        pass
    match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", text, flags=re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except Exception:
            pass
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(text[start : end + 1])
        except Exception:
            pass
    return None


def as_list(value: Any, max_items: int = 5) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        out = []
        for v in value:
            if isinstance(v, dict):
                text = v.get("text") or v.get("description") or " — ".join([str(x) for x in v.values() if x])
                out.append(str(text).strip())
            elif str(v).strip():
                out.append(str(v).strip())
        return out[:max_items]
    lines = [line.strip(" -•\t") for line in str(value).split("\n") if line.strip()]
    return lines[:max_items]


def df_to_bullets(df: pd.DataFrame, cols: List[str], max_items: int = 4) -> List[str]:
    if df.empty:
        return []
    bullets = []
    for _, row in df.head(max_items).iterrows():
        parts = [str(row.get(col, "")).strip() for col in cols if str(row.get(col, "")).strip()]
        bullets.append(" — ".join(parts))
    return bullets


def default_profile_sections(company: pd.Series, project: pd.Series, intel: Dict[str, pd.DataFrame], latest_articles: List[Dict[str, str]]) -> Dict[str, Any]:
    c = company.to_dict()
    financials = intel.get("financials", pd.DataFrame())
    f_row = financials.iloc[0].to_dict() if not financials.empty else {}
    milestones = []
    for _, row in intel.get("milestones", pd.DataFrame()).head(5).iterrows():
        milestones.append({"year": str(row.get("milestone_year", "")), "text": str(row.get("milestone_text", ""))})

    latest_news = []
    for article in latest_articles[:5]:
        title = article.get("title", "").strip()
        source = article.get("source", "").strip()
        if title:
            suffix = f" — {source}" if source else ""
            latest_news.append(f"{title}{suffix} — verify relevance before external use")
    if not latest_news:
        latest_news = ["Latest external news was excluded or unavailable.", "Use internal records and analyst notes as the primary evidence."]

    no_internal_match = str(c.get("company_id", "")) == "external_query"
    verification_banner = "Not found in the internal database — this profile uses external/contextual information and needs additional verification." if no_internal_match else ""

    return {
        "headline": f"One Page Company Profile: {c.get('company_name', 'Company')}",
        "verification_banner": verification_banner,
        "company_positioning": str(c.get("mission", "Mission not available.")),
        "growth_direction": str(c.get("expansion_signal", "Growth direction needs to be verified.")),
        "target_market": str(c.get("target_market", project.get("target_sector", "Target market to be confirmed."))),
        "company_description": str(c.get("short_description", f"{c.get('company_name', 'The company')} operates in {c.get('sector', 'its sector')}." )).strip(),
        "company_snapshot": {
            "hq": str(c.get("hq_location", "Not available")),
            "founded": str(c.get("founded_year", "Not available")),
            "company_type": str(c.get("company_type", "Not available")),
            "sector": str(c.get("sector", "Not available")),
            "employees": str(c.get("employee_count", "Not available")),
        },
        "what_they_do": df_to_bullets(intel.get("products", pd.DataFrame()), ["product_name", "description"], 4) or [str(c.get("recent_news", "Products to verify."))],
        "leadership": df_to_bullets(intel.get("leadership", pd.DataFrame()), ["executive_name", "role"], 3) or ["Leadership not available in internal database."],
        "key_facts": [
            f"HQ: {c.get('hq_location', 'Not available')}",
            f"Founded: {c.get('founded_year', 'Not available')}",
            f"Type: {c.get('company_type', 'Not available')}",
            f"Relevance score: {c.get('relevance_score', 'Not available')}",
        ],
        "funding_commercial_signals": [
            str(f_row.get("latest_funding_or_valuation", "Funding/commercial data not available.")),
            str(f_row.get("commercial_signals", "Commercial signals require verification.")),
            str(f_row.get("financial_caveat", "Private-company financials should be treated cautiously.")),
        ],
        "latest_news_signals": latest_news,
        "risks": [
            "Private-company financials are incomplete unless the source explicitly discloses them.",
            "Latest RSS/news results may be duplicated, irrelevant or behind paywalls.",
            "Treat this as a first draft requiring human review before external use.",
        ],
        "timeline": milestones or [{"year": "TBC", "text": "Milestones not available in internal database."}],
        "next_steps": [
            "Verify leadership, funding and product details against current official sources.",
            "Check whether latest news is relevant and material.",
            "Refine the PowerPoint before sharing with stakeholders.",
        ],
    }


def normalise_profile_sections(raw_text: str, base: Dict[str, Any]) -> Dict[str, Any]:
    parsed = extract_json(raw_text) if raw_text else None
    if not parsed:
        return base
    original_latest_news = as_list(base.get("latest_news_signals"), max_items=5)
    original_timeline = base.get("timeline", [])
    for key in base:
        if key in parsed and parsed[key]:
            base[key] = parsed[key]

    # Keep the latest-news section tied to the first five raw RSS feed items.
    # This avoids empty rows and prevents the model from replacing live RSS
    # titles with older pretrained knowledge.
    if original_latest_news:
        base["latest_news_signals"] = original_latest_news[:5]

    # If the model returns no useful timeline, keep the fallback timeline so the
    # milestone chart is never empty.
    if not base.get("timeline") and original_timeline:
        base["timeline"] = original_timeline
    return base


def profile_to_markdown(profile: Dict[str, Any]) -> str:
    def bullets(key: str, max_items: int = 5) -> str:
        return "\n".join([f"- {x}" for x in as_list(profile.get(key), max_items=max_items)])

    timeline_lines = []
    timeline = profile.get("timeline", [])
    if isinstance(timeline, list):
        for item in timeline[:5]:
            if isinstance(item, dict):
                timeline_lines.append(f"- **{item.get('year', '')}:** {item.get('text', '')}")
            else:
                timeline_lines.append(f"- {item}")

    return f"""
# {profile.get('headline', 'One Page Company Profile')}

{('> **Verification note:** ' + profile.get('verification_banner', '') + '\n') if profile.get('verification_banner') else ''}

## Mission & positioning
{profile.get('company_positioning', '')}

## Growth direction
{profile.get('growth_direction', '')}

## Target market
{profile.get('target_market', '')}

## Company description
{profile.get('company_description', '')}

## What they do
{bullets('what_they_do')}

## Executive leadership
{bullets('leadership')}

## Key facts
{bullets('key_facts')}

## Funding & commercial signals
{bullets('funding_commercial_signals')}

## Latest news / recent signals
{bullets('latest_news_signals')}

## Risks / things to verify
{bullets('risks')}

## Timeline / milestones
{chr(10).join(timeline_lines)}

## Recommended next steps
{bullets('next_steps')}
""".strip()


# -----------------------------
# PowerPoint generation layer
# -----------------------------
# The app now uses the uploaded PowerPoint as the actual output template.
# Keep this file in the repository root when deploying to Streamlit Cloud.
def resolve_template_path() -> Optional[Path]:
    candidates = [
        BASE / "one_pager_template.pptx",
        BASE / "templates" / "one_pager_template.pptx",
        BASE / "Anthropic_One_Pager_v2 (1).pptx",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def truncate_text(text: Any, max_chars: Optional[int] = 180, add_ellipsis: bool = False) -> str:
    value = str(text or "").strip()
    value = re.sub(r"\s+", " ", value)
    if not max_chars or len(value) <= max_chars:
        return value
    suffix = "…" if add_ellipsis else ""
    cut_len = max_chars - len(suffix) if add_ellipsis else max_chars
    return value[:cut_len].rstrip() + suffix


def clean_bullets(value: Any, max_items: int = 5, max_chars: Optional[int] = None, add_ellipsis: bool = False) -> List[str]:
    return [truncate_text(item, max_chars=max_chars, add_ellipsis=add_ellipsis) for item in as_list(value, max_items=max_items)]


def split_name_role(value: str) -> tuple[str, str]:
    value = str(value or "").strip()
    for sep in [" — ", " – ", " - ", ";"]:
        if sep in value:
            left, right = value.split(sep, 1)
            return truncate_text(left, 42), truncate_text(right, 56)
    return truncate_text(value, 42), "Role / relevance to verify"


def set_shape_text(
    slide,
    idx: int,
    text: Any,
    font_size: Optional[float] = None,
    bold: Optional[bool] = None,
    align=None,
    max_chars: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
):
    """Replace text in an existing template shape while keeping full, complete text.

    The previous version sometimes shortened strings before they reached PowerPoint.
    This version only applies max_chars when a caller explicitly requests it for tiny
    fields such as score/year. For body text, it relies on small fonts, wrapping,
    lower margins and resized template boxes so full sentences remain visible.
    """
    if idx >= len(slide.shapes):
        return
    shape = slide.shapes[idx]
    if not hasattr(shape, "text_frame"):
        return
    value = truncate_text(text, max_chars=max_chars) if max_chars else str(text or "")
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    # Minimise internal padding so text is less likely to clip in compact template boxes.
    for attr in ["margin_left", "margin_right", "margin_top", "margin_bottom"]:
        try:
            setattr(tf, attr, Pt(0.5))
        except Exception:
            pass
    lines = str(value).splitlines() or [""]
    for line_no, line in enumerate(lines):
        p = tf.paragraphs[0] if line_no == 0 else tf.add_paragraph()
        p.text = line
        if font_size is not None:
            p.font.size = Pt(font_size)
        if bold is not None:
            p.font.bold = bold
        if font_color is not None:
            p.font.color.rgb = font_color
        if align is not None:
            p.alignment = align
        try:
            p.space_after = Pt(0)
            p.space_before = Pt(0)
            p.line_spacing = 0.92
        except Exception:
            pass


def resize_shape(slide, idx: int, left=None, top=None, width=None, height=None):
    """Resize/reposition a template text box in inches to prevent clipping."""
    if idx >= len(slide.shapes):
        return
    shape = slide.shapes[idx]
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)



def format_top_card_text(value: Any, max_words: int = 26) -> str:
    """Prepare copy for the three top narrative cards without changing template geometry.

    The original PowerPoint template already contains correctly-sized body text
    boxes for Mission / Growth / Target Market. This function therefore returns
    one compact complete sentence and lets PowerPoint wrap it naturally inside
    the existing text box. It does not insert manual line breaks, resize the box,
    or add ellipses.
    """
    if isinstance(value, list):
        raw = " ".join(str(v).strip() for v in value if str(v).strip())
    else:
        raw = str(value or "").strip()
    raw = re.sub(r"\s+", " ", raw)
    if not raw:
        return "To verify."

    words = raw.split()
    if len(words) > max_words:
        raw = " ".join(words[:max_words]).rstrip(" ,;:")

    if not raw.endswith((".", "!", "?")):
        raw += "."
    return raw

def build_slide_values(profile: Dict[str, Any], company: pd.Series, brief_type: str) -> Dict[str, Any]:
    c = company.to_dict() if hasattr(company, "to_dict") else dict(company)
    company_name = str(c.get("company_name", "Company") or "Company")
    snapshot = profile.get("company_snapshot", {}) if isinstance(profile.get("company_snapshot", {}), dict) else {}

    def prefer_profile_snapshot(snapshot_key: str, company_key: str, fallback: str) -> str:
        profile_value = str(snapshot.get(snapshot_key, "") or "").strip()
        company_value = str(c.get(company_key, "") or "").strip()
        unavailable_terms = {"", "not available", "not available from internal database", "nan", "none"}
        if profile_value.lower() not in unavailable_terms:
            return profile_value
        if company_value.lower() not in unavailable_terms:
            return company_value
        return fallback

    sector = prefer_profile_snapshot("sector", "sector", "Sector to verify")
    company_type = prefer_profile_snapshot("company_type", "company_type", "Type to verify")
    relevance_score = str(c.get("relevance_score", "N/A") or "N/A")

    verification_banner = str(profile.get("verification_banner", "") or "").strip()
    if str(c.get("company_id", "")) == "external_query" and not verification_banner:
        verification_banner = "Not found in the internal database — this profile uses external/contextual information and needs additional verification."

    return {
        "company_name": company_name,
        "verification_banner": verification_banner,
        "subtitle": (f"{brief_type}  ·  {sector}" if not verification_banner else f"{brief_type}  ·  External profile — verify"),
        "company_type": company_type,
        "score": relevance_score,
        "mission": format_top_card_text(profile.get("company_positioning", "Mission / positioning to verify.")),
        "growth": format_top_card_text(profile.get("growth_direction", "Growth direction to verify.")),
        "target_market": format_top_card_text(profile.get("target_market", "Target market to verify.")),
        "hq": prefer_profile_snapshot("hq", "hq_location", "HQ to verify"),
        "founded": prefer_profile_snapshot("founded", "founded_year", "Founded year to verify"),
        "type": company_type,
        "sector": sector,
        "employees": prefer_profile_snapshot("employees", "employee_count", "Employees to verify"),
        "description": ((verification_banner + " ") if verification_banner else "") + str(profile.get("company_description", "Company description to verify.")),
        "what_they_do": clean_bullets(profile.get("what_they_do"), max_items=3, max_chars=None),
        "leadership": clean_bullets(profile.get("leadership"), max_items=3, max_chars=None),
        "funding": clean_bullets(profile.get("funding_commercial_signals"), max_items=3, max_chars=None),
        "news": clean_bullets(profile.get("latest_news_signals"), max_items=5, max_chars=None),
        "risks": clean_bullets(profile.get("risks"), max_items=3, max_chars=None),
        "next_steps": clean_bullets(profile.get("next_steps"), max_items=3, max_chars=None),
        "timeline": profile.get("timeline", []),
    }


def add_profile_pptx(profile: Dict[str, Any], company: pd.Series, project: pd.Series) -> BytesIO:
    """Populate the provided one-page PowerPoint template for any company.

    The template is a single-slide company profile. The app keeps the design,
    shapes, icons and layout from the template, and only replaces the text fields
    with AI-generated structured content.
    """
    template_path = resolve_template_path()
    if template_path is None:
        raise FileNotFoundError(
            "PowerPoint template not found. Upload one_pager_template.pptx to the same GitHub folder as app.py."
        )

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    brief_type = str(project.get("project_name", "Company brief")) if hasattr(project, "get") else "Company brief"
    v = build_slide_values(profile, company, brief_type)
    white = RGBColor(255, 255, 255)
    top_body_blue = RGBColor(185, 200, 230)

    # Resize the most content-heavy template text boxes so full sentences can fit.
    # This keeps the visual style of the supplied template but gives body text more room.
    # Top narrative cards: keep the original template box sizes.
    # Left description panel
    resize_shape(slide, 42, height=1.35)
    # Three feature columns
    for idx in [48, 50, 52, 78, 80, 82]:
        resize_shape(slide, idx, height=0.48)
    # Leadership roles
    for idx in [60, 64, 68]:
        resize_shape(slide, idx, height=0.22)
    # News rows: keep five articles, but use the full width of the news panel to avoid clipping.
    news_rows = [88, 90, 92, 94, 96]
    news_tops = [4.46, 4.64, 4.82, 5.00, 5.18]
    for idx, top in zip(news_rows, news_tops):
        resize_shape(slide, idx, left=3.54, top=top, width=9.20, height=0.20)
    # Risks / next steps rows
    for idx in [102, 104, 106, 112, 114, 116]:
        resize_shape(slide, idx, height=0.28)
    # Milestone labels: make each timeline label wider and taller.
    milestone_lefts = [1.65, 4.08, 6.52, 8.96, 11.02]
    for idx, left in zip([122, 125, 128, 131, 134], milestone_lefts):
        resize_shape(slide, idx, left=left, width=1.85, height=0.34)

    # Header / title area — keep the template's original header colours.
    set_shape_text(slide, 2, v["company_name"], font_size=24, bold=True, max_chars=44)
    set_shape_text(slide, 3, v["subtitle"], font_size=9, max_chars=82)
    set_shape_text(slide, 5, v["company_type"], font_size=8.5, align=PP_ALIGN.CENTER, max_chars=42)
    set_shape_text(slide, 7, v["score"], font_size=20, bold=True, align=PP_ALIGN.CENTER, max_chars=5)
    set_shape_text(slide, 8, "SCORE", font_size=6.5, bold=True, align=PP_ALIGN.CENTER)

    # Top three narrative cards — keep the template's original text-box sizes.
    # Only replace the body copy; do not resize or reposition these text boxes.
    # The original template headings remain in place.
    set_shape_text(slide, 13, v["mission"], font_size=7.2, font_color=top_body_blue)
    set_shape_text(slide, 18, v["growth"], font_size=7.2, font_color=top_body_blue)
    set_shape_text(slide, 23, v["target_market"], font_size=7.2, font_color=top_body_blue)

    # Left company snapshot panel — these values sit on the dark left panel, so use white text.
    set_shape_text(slide, 27, v["hq"], font_size=7.6, max_chars=42, font_color=white)
    set_shape_text(slide, 30, v["founded"], font_size=7.6, max_chars=18, font_color=white)
    set_shape_text(slide, 33, v["type"], font_size=7.6, max_chars=42, font_color=white)
    set_shape_text(slide, 36, v["sector"], font_size=7.6, max_chars=46, font_color=white)
    set_shape_text(slide, 39, v["employees"], font_size=7.6, max_chars=42, font_color=white)
    set_shape_text(slide, 42, v["description"], font_size=4.9, font_color=white)

    # What they do
    for idx, text in zip([48, 50, 52], v["what_they_do"] + ["Not available / to verify."] * 3):
        set_shape_text(slide, idx, text, font_size=5.0)

    # Leadership — only three current global leaders are shown on the slide.
    leadership = v["leadership"] + ["Leadership data to verify — Role to verify"] * 3
    for item, name_idx, role_idx in zip(leadership[:3], [59, 63, 67], [60, 64, 68]):
        name, role = split_name_role(item)
        set_shape_text(slide, name_idx, name, font_size=7.0, bold=True, max_chars=42)
        set_shape_text(slide, role_idx, role, font_size=4.8)
    # Clear the fourth leadership slot in the template.
    set_shape_text(slide, 71, "", font_size=7.0)
    set_shape_text(slide, 72, "", font_size=5.8)

    # Funding / commercial signals
    for idx, text in zip([78, 80, 82], v["funding"] + ["Funding / commercial signal to verify."] * 3):
        set_shape_text(slide, idx, text, font_size=5.0)

    # Latest news / signals — the template has five rows across the wide panel
    set_shape_text(slide, 86, "LATEST NEWS / SIGNALS  ·  VERIFY RELEVANCE BEFORE EXTERNAL USE", font_size=8.6, bold=True, max_chars=80)
    for idx, text in zip([88, 90, 92, 94, 96], v["news"] + ["No recent article used / verify external signals."] * 5):
        # Do not truncate article titles with ellipses; let PowerPoint shrink/wrap text to fit.
        set_shape_text(slide, idx, text, font_size=4.1, max_chars=None)

    # Risks and next steps
    for idx, text in zip([102, 104, 106], v["risks"] + ["Risk / caveat to verify."] * 3):
        set_shape_text(slide, idx, text, font_size=4.9)
    for idx, text in zip([112, 114, 116], v["next_steps"] + ["Next step to verify."] * 3):
        set_shape_text(slide, idx, text, font_size=4.9)

    # Timeline / milestones
    timeline = v["timeline"] if isinstance(v["timeline"], list) else []
    if not timeline:
        timeline = [{"year": "TBC", "text": "Milestones to verify"}]
    timeline = timeline[:5]
    year_shapes = [121, 124, 127, 130, 133]
    text_shapes = [122, 125, 128, 131, 134]
    for i, (year_idx, text_idx) in enumerate(zip(year_shapes, text_shapes)):
        item = timeline[i] if i < len(timeline) else {"year": "", "text": ""}
        if not isinstance(item, dict):
            item = {"year": "", "text": str(item)}
        set_shape_text(slide, year_idx, item.get("year", ""), font_size=7.2, bold=True, align=PP_ALIGN.CENTER, max_chars=8)
        set_shape_text(slide, text_idx, item.get("text", ""), font_size=3.7, align=PP_ALIGN.CENTER, max_chars=None)

    footer_text = "AI-generated first draft — funding, leadership and news signals require human verification before external use."
    if v.get("verification_banner"):
        footer_text = v["verification_banner"] + " Human review required before use."
    set_shape_text(
        slide,
        135,
        footer_text,
        font_size=5.9,
        align=PP_ALIGN.RIGHT,
    )

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# -----------------------------
# Streamlit UI
# -----------------------------
st.title("BriefIQ")
st.caption("Prototype: free-text company request + internal company intelligence database + optional latest news + LLM workflow → editable PowerPoint one-pager")
st.caption("Data freshness note: internal database records are static demo data; the LLM has a fixed training cutoff, so current leadership/funding facts should be verified against database, user-provided context or latest news before use.")

if not DB_PATH.exists():
    try:
        import setup_demo_data
        setup_demo_data.main()
        DB_PATH = resolve_db_path()
        st.info("Demo database was missing, so it has been created automatically.")
    except Exception as e:
        st.error(f"Demo database could not be created automatically: {e}")
        st.stop()

companies = safe_load_table("companies", str(DB_PATH))
if companies.empty:
    st.error("No companies table found in the demo database. Re-upload briefing_demo.db or setup_demo_data.py.")
    st.stop()
notes = read_markdown_notes()

with st.sidebar:
    st.header("Demo controls")
    provider = st.selectbox("LLM mode", ["OpenAI", "Demo fallback"])
    if provider == "OpenAI" and not get_secret("OPENAI_API_KEY"):
        st.warning("OpenAI key not found. Add OPENAI_API_KEY in Streamlit secrets, or use Demo fallback.")
    st.divider()
    st.caption(f"Database path: {DB_PATH.relative_to(BASE) if DB_PATH.exists() else DB_PATH}")
    st.caption(f"Internal company records: {len(companies)}")
    st.caption(f"Markdown notes found: {len(notes)}")
    for table in ["leadership", "products", "financials", "milestones"]:
        st.caption(f"{table.title()} records: {len(safe_load_table(table, str(DB_PATH)))}")

st.subheader("1. Enter company and brief type")
company_query = st.text_input("Company name", value="Anthropic", placeholder="Type any company name, for example Anthropic, OpenAI or Mistral AI")
brief_type = st.selectbox("Brief type", ["Investment one-pager", "Company brief"])

st.subheader("2. Add context")
extra_notes = st.text_area(
    "Additional context, messy notes or instructions",
    height=150,
    placeholder="Example: Focus on enterprise AI adoption, leadership, funding/commercial signals, risks, and why this company matters for a senior stakeholder.",
)
uploaded_file = st.file_uploader("Optional: upload a .txt or .md file", type=["txt", "md"])
if uploaded_file is not None:
    extra_notes += "\n\nUPLOADED FILE CONTENT:\n" + uploaded_file.read().decode("utf-8", errors="ignore")

st.subheader("3. Match to internal database")
selected_company, has_internal_match, match_message = find_company_match(company_query, companies)
if has_internal_match:
    st.success(match_message)
else:
    st.warning(match_message)

intel = company_intelligence(selected_company) if has_internal_match else {"leadership": pd.DataFrame(), "products": pd.DataFrame(), "financials": pd.DataFrame(), "milestones": pd.DataFrame()}
with st.expander("View matched internal company intelligence"):
    st.write("Company profile")
    st.dataframe(pd.DataFrame(selected_company).rename(columns={0: "value"}))
    for label, df in intel.items():
        st.write(label.title())
        if df.empty:
            st.caption("No internal records found.")
        else:
            st.dataframe(df)

st.subheader("4. Latest news signals")
include_latest_news = st.checkbox(
    "Include latest 5 news/RSS articles in the one-pager",
    value=True,
    help="If selected, the app retrieves recent Google News RSS results for the company and passes them into the LLM as external signals to verify.",
)
if include_latest_news:
    st.info("Latest news will be retrieved when you generate the one-pager. It is a signal source, not definitive evidence.")
else:
    st.info("Latest news is excluded. The one-pager will rely on internal data, notes and your added context.")

st.subheader("5. Generate output")
st.write("This runs the workflow: interpret request → match internal company intelligence → retrieve notes → optional news lookup → live leadership lookup where needed → draft structured profile → review → export editable PowerPoint.")

if st.button("Generate one-pager", type="primary"):
    selected_project = build_project_context(brief_type, extra_notes)
    query = f"{company_query} {selected_company.get('company_name', '')} {brief_type} {extra_notes}"
    retrieved_notes = retrieve_relevant_notes(query, notes)

    with st.spinner("Retrieving latest news signals..." if include_latest_news else "Preparing source pack..."):
        latest_articles = get_latest_company_articles(selected_company.get("company_name", company_query), max_articles=5) if include_latest_news else []
        internal_leadership_missing = intel.get("leadership", pd.DataFrame()).empty
        live_leadership_lookup = ""
        live_company_profile_lookup = ""
        if provider == "OpenAI" and not has_internal_match:
            live_company_profile_lookup = lookup_current_company_profile(str(selected_company.get("company_name", company_query)))
        if provider == "OpenAI" and internal_leadership_missing:
            live_leadership_lookup = lookup_current_global_leadership(str(selected_company.get("company_name", company_query)))
        source_pack = build_source_pack(
            selected_company,
            selected_project,
            extra_notes,
            retrieved_notes,
            latest_articles,
            intel,
            has_internal_match=has_internal_match,
            match_message=match_message,
            live_leadership_lookup=live_leadership_lookup,
            live_company_profile_lookup=live_company_profile_lookup,
        )

    base_profile = default_profile_sections(selected_company, selected_project, intel, latest_articles)
    with st.spinner("Drafting structured profile..."):
        if provider == "Demo fallback":
            raw_draft = json.dumps(base_profile, indent=2)
            profile = base_profile
        else:
            draft_prompt = build_prompt(source_pack, brief_type)
            raw_draft = run_llm(draft_prompt, provider)
            profile = normalise_profile_sections(raw_draft, base_profile)
        markdown_draft = profile_to_markdown(profile)

    with st.spinner("Running review step..."):
        if provider == "Demo fallback":
            review = "Demo review: verify all leadership, funding/valuation and latest-news claims before external use. Treat private-company financial data as directional commercial signals, not full financial performance."
        else:
            review_prompt = build_review_prompt(json.dumps(profile, indent=2), source_pack)
            review = run_llm(review_prompt, provider)

    pptx_file = add_profile_pptx(profile, selected_company, selected_project)
    clean_name = re.sub(r"[^A-Za-z0-9_]+", "_", str(selected_company.get("company_name", company_query))).strip("_") or "company"
    st.session_state["raw_draft"] = raw_draft
    st.session_state["profile"] = profile
    st.session_state["draft"] = markdown_draft
    st.session_state["review"] = review
    st.session_state["source_pack"] = source_pack
    st.session_state["latest_articles"] = latest_articles
    st.session_state["include_latest_news"] = include_latest_news
    st.session_state["match_message"] = match_message
    st.session_state["pptx_file"] = pptx_file.getvalue()
    st.session_state["pptx_name"] = f"{clean_name}_{brief_type.replace(' ', '_').lower()}.pptx"

if "profile" in st.session_state:
    tab1, tab2, tab3, tab4, tab5 = st.tabs(["PowerPoint output", "Text preview", "Latest news", "AI review", "Source pack"])
    with tab1:
        st.success("PowerPoint one-pager created. Download and open in PowerPoint to edit the slide.")
        st.caption(st.session_state.get("match_message", ""))
        st.download_button(
            "Download one-page profile as PowerPoint",
            data=st.session_state["pptx_file"],
            file_name=st.session_state.get("pptx_name", "one_page_profile.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        with st.expander("View structured fields used in the slide"):
            st.json(st.session_state["profile"])
    with tab2:
        st.markdown(st.session_state["draft"])
        st.download_button("Download text preview as Markdown", data=st.session_state["draft"], file_name="ai_one_pager_preview.md", mime="text/markdown")
    with tab3:
        if not st.session_state.get("include_latest_news"):
            st.write("Latest news was excluded for this run.")
        elif st.session_state.get("latest_articles"):
            for article in st.session_state.get("latest_articles", []):
                st.markdown(f"**{article.get('title', '')}**")
                if article.get("published"):
                    st.caption(article.get("published"))
                if article.get("summary"):
                    st.write(article.get("summary"))
                if article.get("link"):
                    st.write(article.get("link"))
                st.divider()
        else:
            st.write("No latest articles retrieved for this company.")
    with tab4:
        st.markdown(st.session_state["review"])
    with tab5:
        st.code(st.session_state["source_pack"], language="text")

st.divider()
st.caption(
    "Interview talking point: this version starts with a free-text company request, checks a richer internal company intelligence database "
    "for profile, mission, leadership, products, funding/commercial signals and milestones, optionally retrieves latest external signals, "
    "then uses an LLM workflow to structure the evidence into an editable PowerPoint one-pager."
)
